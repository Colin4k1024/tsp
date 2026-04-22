#!/usr/bin/env python3
"""Generate Team Skills Platform Introduction PPTX"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Content
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.space_after = Pt(12)

    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(0.5))
    lt = left_title_box.text_frame.paragraphs[0]
    lt.text = left_title
    lt.font.size = Pt(22)
    lt.font.bold = True

    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.3), Inches(5))
    lt2 = left_box.text_frame
    lt2.word_wrap = True
    for i, b in enumerate(left_bullets):
        p = lt2.paragraphs[0] if i == 0 else lt2.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(16)
        p.space_after = Pt(8)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(0.5))
    rt = right_title_box.text_frame.paragraphs[0]
    rt.text = right_title
    rt.font.size = Pt(22)
    rt.font.bold = True

    # Right content
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.3), Inches(5))
    rt2 = right_box.text_frame
    rt2.word_wrap = True
    for i, b in enumerate(right_bullets):
        p = rt2.paragraphs[0] if i == 0 else rt2.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(16)
        p.space_after = Pt(8)

    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    add_title_slide(
        prs,
        "Team Skills Platform",
        "角色化研发团队协作平台介绍 | 2026-03-29"
    )

    # Slide 2: Problem Background
    add_content_slide(prs, "问题背景：现状痛点", [
        "单 Agent 能力局限：单个 AI 无法承担完整研发链路",
        "角色边界模糊：缺乏清晰的职责分工和交接契约",
        "工程规范碎片化：前端、API、安全、测试规范分散不统一",
        "平台绑定：规范定义无法同时适配 Codex 与 Claude 双平台",
        "技能传承困难：最佳实践依赖个人经验，难以系统化沉淀"
    ])

    # Slide 3: Solution
    add_content_slide(prs, "解决方案：核心思路", [
        "8 大角色矩阵：Tech Lead 统筹 + 7 个专业角色分工协作",
        "标准化交接契约：角色间交接有明确定义和 quality gates",
        "Canonical Source of Truth：roles/ 目录为单一真相源",
        "双平台自动生成：Codex 与 Claude 技能/Agent 自动导出",
        "ECC Harness Layer：27 个专家级 Specialist Agent 增强"
    ])

    # Slide 4: Architecture Overview
    add_two_column_slide(
        prs,
        "架构概览",
        "8 角色矩阵",
        [
            "tech-lead: 统一 intake、拆解、分派、冲突决策",
            "product-manager: 需求、PRD、用户故事、验收标准",
            "project-manager: 排期、依赖、里程碑、风险",
            "architect: ADR、接口契约、数据边界",
            "frontend-engineer: 前端实现与自测",
            "backend-engineer: 后端实现与自测",
            "qa-engineer: 测试计划、回归验证、放行建议",
            "devops-engineer: 发布、监控、回滚与运行保障"
        ],
        "ECC Harness Layer",
        [
            "27 个 Specialist Agents",
            "planner, architect, chief-of-staff",
            "security-reviewer (OWASP Top 10)",
            "build-error-resolver (8 语言版本)",
            "tdd-guide, loop-operator, harness-optimizer",
            "e2e-runner, refactor-cleaner",
            "语言专项审查 (Java/Go/Python/Rust...)"
        ]
    )

    # Slide 5: Engineering Results - Role System
    add_content_slide(prs, "工程化成果：角色与规范体系", [
        "roles/: 8 个角色的 canonical 定义（role.yaml）",
        "skills/roles/: 自动生成的角色技能（平台无关）",
        "agents/roles/: 自动生成的角色 Agent 提示词",
        "agents/specialists/: 27 个 ECC 专家级 Agent",
        "commands/: 团队命令表面（/team-intake, /team-plan, /handoff 等）",
        "rules/: 团队工作规则 + 语言规范包（TS/Java/Python/Go）"
    ])

    # Slide 6: Engineering Results - Frontend Governance
    add_two_column_slide(
        prs,
        "工程化成果：前端治理",
        "前端工程规范",
        [
            "frontend-engineering: React/Next 工程规范",
            "组件结构、状态管理、语义化",
            "可访问性（Accessibility）基线",
            "性能优化标准"
        ],
        "UI/UX 设计系统",
        [
            "frontend-ui-ux-system: 设计知识库",
            "产品类型、视觉方向、设计 Token",
            "交互模式、响应式设计",
            "前端质量门禁规则"
        ]
    )

    # Slide 7: Engineering Results - Security Gates
    add_content_slide(prs, "工程化成果：安全与质量门禁", [
        "CodeQL 安全扫描：PR 级别代码安全审查",
        "Secret 扫描：防止密钥/凭证泄露",
        "SBOM 生成：软件物料清单追溯",
        "SLSA 验证：供应链安全级别认证",
        "Contract Testing：API 契约测试",
        "Trivy/Checkov：容器与 IaC 安全扫描",
        "Security Reviewer：OWASP Top 10 全覆盖（30+ 检查项）"
    ])

    # Slide 8: Hooks System
    add_two_column_slide(
        prs,
        "工程化成果：Hooks 系统",
        "PreToolUse Hooks (5个)",
        [
            "tmux-reminder — dev server 建议用 tmux",
            "git-push-reminder — push 前检查变更",
            "console-log-warning — 检测调试语句",
            "doc-file-warning — 非标准文档警告",
            "mcp-health-check — MCP 健康检查"
        ],
        "PostToolUse Hooks (3个)",
        [
            "quality-gate — 质量检查（空格/分号）",
            "format-suggestion — JS/TS 格式化建议",
            "pr-created — PR URL 自动提取"
        ]
    )

    # Slide 9: Continuous Learning & Cost Awareness
    add_two_column_slide(
        prs,
        "工程化成果：持续学习 + 成本感知",
        "Continuous Learning v2.1",
        [
            "Instinct Store：项目级/全局级 instinct",
            "observe hook：工具调用观察捕获",
            "instinct_cli：list/evolve/export/import",
            "自动从会话提取最佳实践模式",
            "置信度评分（0.3-0.9）驱动演进"
        ],
        "Cost-Aware LLM Pipeline",
        [
            "Haiku：轻量任务（$0.001/1K）",
            "Sonnet：主开发（$0.01/1K）",
            "Opus：深度推理（$0.05/1K）",
            "Context Budget：70%/85% 阈值预警",
            "cost_tracker：会话成本记录"
        ]
    )

    # Slide 10: Public Commands Expansion
    add_two_column_slide(
        prs,
        "工程化成果：公开命令面扩展",
        "主链 + specialist",
        [
            "/team-intake -> /team-plan -> /team-execute",
            "/handoff -> /team-review -> /team-release",
            "/plan、/code-review、/build-fix、/verify",
            "/multi-frontend、/multi-backend"
        ],
        "本轮新增重点",
        [
            "/tdd — 先锁测试与成功标准",
            "/harness-audit — 平台能力体检与缺口收敛",
            "specialist 结果必须回收到主链",
            "平台治理与业务交付路径分离"
        ]
    )

    # Slide 11: Runtime vs Explicit Commands
    add_two_column_slide(
        prs,
        "工程化成果：显式命令与 Runtime 分层",
        "显式命令",
        [
            "用户可见入口，负责触发决策与流程",
            "例：/team-plan、/tdd、/harness-audit",
            "适合解释主链、specialist 与角色协作",
            "对外演示时最容易理解"
        ],
        "Runtime 后台机制",
        [
            "memory persistence、observe、cost tracking",
            "context budget、compact/archive、instinct learning",
            "自动增强长会话稳定性与成本控制",
            "不是每次都手动执行的命令"
        ]
    )

    # Slide 12: Core Advantages
    add_content_slide(prs, "核心优势", [
        "Canonical Source of Truth：单一真相源，平台无关定义",
        "双平台支持：Codex 与 Claude 同步覆盖",
        "角色化协作：明确职责边界，标准化交接",
        "ECC 增强层：27 个专家级 Agent 快速编排",
        "Hooks 系统：8 个自动化质量门禁钩子",
        "持续学习：Instinct-based 自动模式提取",
        "120+ Runbooks：覆盖快速入门、使用场景、治理、故障排查",
        "可安装、可验证、可迁移：不是项目脚手架，而是能力基座"
    ])

    # Slide 13: Demo Materials
    add_content_slide(prs, "演示与汇报材料", [
        "业务主链演示：demo-scenario.md + demo-execution-log.md",
        "平台升级演示：platform-capability-demo-script.md + execution-log",
        "vertical 场景导览：vertical-scenario-route-map.md",
        "vertical 场景矩阵：vertical-scenario-capability-matrix.md",
        "已具备完整 demo 闭环的 9 类 vertical 场景可直接复用",
        "命令总表：command-and-capability-matrix.md",
        "runtime 总览：runtime-capabilities-overview.md",
        "对外讲解可直接配合 presentation-talk-track.md"
    ])

    # Slide 14: Vertical Demo Route Map
    add_two_column_slide(
        prs,
        "vertical Demo 成熟度",
        "已具备完整闭环",
        [
            "GitHub Actions / 供应链治理",
            "AI / Eval 平台",
            "移动端 / 小程序交付",
            "IaC / Kubernetes 平台仓库",
            "插件 / 扩展仓库",
            "数据 / ML pipeline 仓库",
            "安全 / 合规平台仓库",
            "内部开发者平台",
            "数据可观测性 / 质量平台"
        ],
        "推荐使用方式",
        [
            "onboarding：先看 route map，再用 matrix 查缺层",
            "汇报：直接复用各 vertical demo script",
            "复盘：搭配 execution log 做团队同步",
            "平台讲解：先讲 platform capability，再切到 vertical"
        ]
    )

    # Slide 15: Next Steps
    add_content_slide(prs, "下一步规划", [
        "平台接入：扩大团队试点范围",
        "规范完善：补充更多语言/框架规范包",
        "ECC 扩展：继续沉淀专家级技能",
        "度量体系：建立团队效能度量指标",
        "知识沉淀：持续丰富 Runbooks 库",
        "Eval 驱动开发：完善 pass@k 指标体系"
    ])

    # Save
    output_path = "/Users/jiafan/Desktop/poc/harness-demo/docs/presentation/team-skills-platform-intro.pptx"
    prs.save(output_path)
    print(f"PPTX saved to: {output_path}")

if __name__ == "__main__":
    main()
