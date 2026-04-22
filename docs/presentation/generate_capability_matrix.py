#!/usr/bin/env python3
"""Generate Agent Capability Matrix PPTX - Updated with new capabilities"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_title(slide, text, top=0.4):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    return txBox

def add_bullets(slide, bullets, left=0.5, top=1.4, width=9, font_size=18, line_spacing=24):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f'• {b}'
        p.font.size = Pt(font_size)
        p.space_after = Pt(line_spacing - font_size)
    return txBox

def add_two_columns(slide, title, left_title, left_items, right_title, right_items, top=0.4):
    set_title(slide, title, top)
    lt_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(0.5))
    lt = lt_box.text_frame.paragraphs[0]
    lt.text = left_title
    lt.font.size = Pt(20)
    lt.font.bold = True
    lb = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.3), Inches(5))
    ltf = lb.text_frame
    ltf.word_wrap = True
    for i, b in enumerate(left_items):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = f'• {b}'
        p.font.size = Pt(16)
        p.space_after = Pt(8)
    rt_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(0.5))
    rt = rt_box.text_frame.paragraphs[0]
    rt.text = right_title
    rt.font.size = Pt(20)
    rt.font.bold = True
    rb = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.3), Inches(5))
    rtf = rb.text_frame
    rtf.word_wrap = True
    for i, b in enumerate(right_items):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = f'• {b}'
        p.font.size = Pt(16)
        p.space_after = Pt(8)

# Slide 1: Cover
s1 = add_slide(prs)
txBox = s1.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = 'Agent 能力矩阵'
p.font.size = Pt(48)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
txBox2 = s1.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(1))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = 'Harness Layer 现状分析与发展规划'
p2.font.size = Pt(24)
p2.alignment = PP_ALIGN.CENTER
txBox3 = s1.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(0.5))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = '2026-03-29'
p3.font.size = Pt(16)
p3.alignment = PP_ALIGN.CENTER

# Slide 2: 能力分类概览
s2 = add_slide(prs)
set_title(s2, 'Agent 能力体系分类')
add_bullets(s2, [
    'Token/Context Management — 上下文窗口管理、模型选择、令牌优化',
    'Memory & Persistence — 跨会话记忆、自动保存/加载、上下文恢复',
    'Learning & Adaptation — 从会话中提取模式、持续优化、错误经验沉淀',
    'Verification & Evaluation — 检查点验证、评分器、pass@k 指标',
    'Parallelization & Scaling — 多实例并行、任务分发、动态扩展',
    'Sub-agent Orchestration — 多代理协作、上下文传递、迭代检索'
], font_size=20, line_spacing=32)

# Slide 3: 你填写的 6 项能力
s3 = add_slide(prs)
set_title(s3, '你已定义的能力范畴')
add_bullets(s3, [
    'Token 优化 — 模型选择、系统提示精简、后台进程',
    '内存持久化 — 自动跨会话保存/加载上下文的钩子',
    '持续学习 — 从会话中自动提取模式到可重用的技能',
    '验证循环 — 检查点 vs 持续评估、评分器类型、pass@k 指标',
    '并行化 — Git worktrees、级联方法、何时扩展实例',
    '子代理编排 — 上下文问题、迭代检索模式'
], font_size=20, line_spacing=32)

# Slide 4: 能力矩阵现状（最新）
s4 = add_slide(prs)
set_title(s4, '能力矩阵现状分析（最新）')
rows = [
    ['能力领域', '当前状态', '实现方式'],
    ['Token/Context Management', '✅ 已实现', 'pre_compact + suggest_compact + trigger_compact'],
    ['Memory & Persistence', '✅ 已实现', 'session_start/end hooks + memory_store'],
    ['Learning & Adaptation', '✅ 已实现', 'instinct_store + observe hook + instinct_cli'],
    ['Verification & Evaluation', '✅ 强覆盖', 'eval-harness skill + pass@k 指标'],
    ['Parallelization & Scaling', '✅ 已实现', 'parallel-execution skill + worktree'],
    ['Sub-agent Orchestration', '✅ 强覆盖', '8 角色矩阵、/team-* 流程链'],
    ['Cost-Aware Pipeline', '✅ 已实现', 'cost_tracker + context_budget + Haiku/Sonnet/Opus'],
    ['Error Handling & Recovery', '✅ 强覆盖', 'build-error-resolver + 7 语言版本'],
    ['Security & Governance', '✅ 强覆盖', 'security-reviewer + OWASP Top 10 + 30+ 检查项'],
    ['PreToolUse Hooks', '✅ 已实现', 'tmux/git/console/doc/mcp-health 五种钩子'],
    ['PostToolUse Hooks', '✅ 已实现', 'quality-gate + format + pr-created 三种钩子'],
]
col_widths = [Inches(2.8), Inches(1.4), Inches(5.3)]
row_height = Inches(0.55)
start_top = Inches(1.3)
left = Inches(0.5)
for i, h in enumerate(rows[0]):
    hdr = s4.shapes.add_textbox(left, start_top, col_widths[i], row_height)
    tf = hdr.text_frame
    p = tf.paragraphs[0]
    p.text = h
    p.font.size = Pt(14)
    p.font.bold = True
    left += col_widths[i]
top = start_top + row_height
for row in rows[1:]:
    left = Inches(0.5)
    for i, cell in enumerate(row):
        cell_box = s4.shapes.add_textbox(left, top, col_widths[i], row_height)
        tf = cell_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cell
        p.font.size = Pt(13)
        left += col_widths[i]
    top += row_height

# Slide 5: Memory & Persistence 详情
s5 = add_slide(prs)
add_two_columns(s5, 'Memory & Persistence 已实现',
    'session_start hook',
    [
        '加载上次会话摘要',
        '检索 pending items',
        '检查待确认决策',
        '提供下次会话提示',
        '存储：~/.claude/memory/sessions/'
    ],
    'session_end hook',
    [
        '收集会话数据',
        '保存任务列表',
        '持久化决策快照',
        '记录 key_findings',
        '自动生成 session_id'
    ]
)

# Slide 6: Context Compression 详情
s6 = add_slide(prs)
set_title(s6, 'Context Compression 已实现')
add_bullets(s6, [
    'pre_compact hook — 压缩前分析，分类 keep/summarize/discard',
    '  └ 高价值：decision、conclusion、pending_item、verification_result',
    '  └ 低价值：tool_result、search_result（>500字符时丢弃）',
    'suggest_compact hook — 上下文使用 >70% 时触发压缩建议',
    '  └ 阈值：low(<70%) / medium(70-85%) / high(85-95%) / critical(>95%)',
    '  └ 4阶段重组计划：保存决策→保存待办→压缩对话→精简工具输出'
], font_size=17, line_spacing=24)

# Slide 7: Error Experience Library
s7 = add_slide(prs)
add_two_columns(s7, 'Learning & Adaptation: Error Experience Library',
    '核心操作',
    [
        'record: 记录错误模式',
        'search: 查询历史模式',
        'feedback: 更新成功率',
        '按 (success - failure) 排序',
        '高成功率优先推荐'
    ],
    '存储结构',
    [
        '~/.claude/memory/error_experience/',
        'patterns/: 错误模式库',
        '  - error_type',
        '  - root_cause + solution',
        '  - success_count / failure_count',
        'decisions/: 决策快照'
    ]
)

# Slide 8: Parallel Execution
s8 = add_slide(prs)
add_two_columns(s8, 'Parallelization & Scaling 已实现',
    '任务类型',
    [
        'independent: 独立任务 → 并行',
        'sequential: 串行依赖',
        'parallel: 多独立任务 → 并行',
        'cascade: 主任务触发子任务'
    ],
    'Git Worktree 并行',
    [
        'worktree add <path> <branch>',
        'worktree list',
        'worktree prune',
        'Scale Out: 增加实例数',
        'Scale Up: 提升实例规格'
    ]
)

# Slide 9: 项目已有强项
s9 = add_slide(prs)
add_two_columns(s9, '项目已有强项',
    'Verification & Evaluation',
    [
        'maven-qa (4阶段质量门禁)',
        'pairwise-test-design',
        'browser-smoke-testing',
        'testcontainers-integration',
        'tdd-guide + e2e-runner',
        '45+ 质量门禁检查项'
    ],
    'Sub-agent Orchestration',
    [
        '8 角色矩阵 (tech-lead 统筹)',
        '/team-* 完整流程链',
        '/handoff 标准化交接',
        'chief-of-staff 跨角色协调',
        '27 个 Specialist Agents',
        '27 个 ECC 技能'
    ]
)

# Slide 10: 新增能力文件清单
s10 = add_slide(prs)
set_title(s10, '新增能力文件清单')
add_bullets(s10, [
    'scripts/lib/memory_store.py — 底层存储接口',
    'scripts/lib/instinct_store.py — Instinct 持久化（项目/全局作用域）',
    'scripts/lib/cost_tracker.py — Token 成本追踪',
    'scripts/lib/context_budget.py — 上下文预算监控',
    'scripts/lib/trigger_compact.py — 触发式压缩（4种触发器）',
    'scripts/lib/context_priority.py — 优先级分类（高/中/低）',
    'scripts/lib/context_archiver.py — 上下文归档到内存',
    'scripts/hooks/pre_compact.py — 增强：trigger-based 压缩',
    'scripts/hooks/cost_tracker.py — Stop 事件成本记录',
    'scripts/hooks/observe.py — 工具调用观察捕获',
    'skills/eval-harness/ — Eval 驱动开发 + pass@k',
    'skills/continuous-learning/ — Instinct 持续学习 v2.1',
    'skills/strategic-compact/ — 上下文优化技能',
    'skills/cost-aware-llm-pipeline/ — 成本感知 LLM 管道',
    'commands/harness-audit.md — 7 维度评分审计命令',
    'docs/runbooks/command-and-capability-matrix.md — 命令与能力矩阵',
    'docs/runbooks/runtime-capabilities-overview.md — runtime 总览',
    'docs/runbooks/platform-capability-demo-script.md — 平台升级演示剧本'
], font_size=14, line_spacing=20)

# Slide 11: PreToolUse/PostToolUse Hooks
s11 = add_slide(prs)
add_two_columns(s11, 'Hooks 系统增强',
    'PreToolUse Hooks (5个)',
    [
        'tmux-reminder — 检测 dev server 建议用 tmux',
        'git-push-reminder — push 前检查未提交变更',
        'console-log-warning — 检测 console.log 等调试语句',
        'doc-file-warning — 警告非标准位置文档',
        'mcp-health-check — MCP 服务器健康检查'
    ],
    'PostToolUse Hooks (3个)',
    [
        'quality-gate — 编辑后质量检查（空格/分号/console）',
        'format-suggestion — JS/TS 文件格式化建议',
        'pr-created — gh pr create 后提取 PR URL'
    ]
)

# Slide 12: Security Reviewer OWASP Top 10
s12 = add_slide(prs)
set_title(s12, 'Security Reviewer OWASP Top 10 覆盖')
add_bullets(s12, [
    'A01: Broken Access Control — 访问控制检查',
    'A02: Cryptographic Failures — 加密失效检查',
    'A03: Injection — 注入漏洞检查',
    'A04: Insecure Design — 不安全设计检查',
    'A05: Security Misconfiguration — 安全配置错误',
    'A06: Vulnerable Components — 漏洞组件检查',
    'A07: Authentication Failures — 认证失败检查',
    'A08: Data Integrity Failures — 数据完整性检查',
    'A09: Logging Failures — 日志记录失败检查',
    'A10: SSRF — 服务端请求伪造检查'
], font_size=17, line_spacing=22)

# Slide 13: 架构演进路线（最新）
s13 = add_slide(prs)
add_two_columns(s13, '架构演进路线（最新）',
    'Phase 1-3: 已完成 ✅',
    [
        'Memory Persistence hooks',
        'Context Compression hooks',
        'Error Experience Library',
        'Parallel Execution Framework',
        'PreToolUse/PostToolUse Hooks',
        'Continuous Learning System',
        'Cost Tracking + Context Budget',
        'Strategic Compact 增强'
    ],
    'Phase 4-7: 已完成 ✅',
    [
        'Eval Harness + pass@k 指标',
        'Cost-Aware LLM Pipeline',
        'Harness Audit 7维度评分',
        'Security Reviewer OWASP Top 10',
        'Instinct-based Learning v2.1',
        'Trigger-based Lazy Loading',
        'context_archiver 上下文归档'
    ]
)

# Slide 14: 新增命令与 runtime 边界
s14 = add_slide(prs)
add_two_columns(s14, '新增命令与 Runtime 边界',
    '显式命令',
    [
        '/tdd — 先定义测试与完成标准',
        '/harness-audit — 平台覆盖度体检',
        '/team-* — 主链角色协作入口',
        'specialist 结果必须回到 handoff 或主链'
    ],
    'Runtime 机制',
    [
        'memory persistence — 会话摘要与待办回存',
        'observe — 工具调用模式观察',
        'cost-aware pipeline — 成本与预算控制',
        'compact/archive — 长会话上下文重组'
    ]
)

# Slide 15: 演示材料与建议阅读路径
s15 = add_slide(prs)
set_title(s15, '演示材料与建议阅读路径')
add_bullets(s15, [
    '首次介绍平台：team-skills-platform-intro.pptx',
    '能力矩阵说明：agent-capability-matrix.pptx',
    '命令全景：command-and-capability-matrix.md',
    'runtime 总览：runtime-capabilities-overview.md',
    '平台升级演示：platform-capability-demo-script.md',
    'vertical 场景导览：docs/presentation/vertical-scenario-route-map.md',
    'vertical 场景矩阵：docs/runbooks/vertical-scenario-capability-matrix.md',
    '讲解提词稿：docs/presentation/presentation-talk-track.md'
], font_size=18, line_spacing=24)

# Slide 16: Vertical 场景成熟度
s16 = add_slide(prs)
add_two_columns(s16, 'vertical 场景成熟度与使用方式',
    '完整 demo 闭环',
    [
        'GitHub Actions / 供应链治理',
        'AI / Eval 平台',
        '移动端 / 小程序交付',
        'IaC / Kubernetes 平台仓库',
        '插件 / 扩展仓库',
        '数据 / ML pipeline 仓库',
        '安全 / 合规平台仓库',
        '内部开发者平台',
        '数据可观测性 / 质量平台'
    ],
    '建议入口',
    [
        '先看 vertical-scenario-route-map.md',
        '再看 vertical-scenario-capability-matrix.md',
        '再按项目类型选对应 demo script',
        '需要复盘时搭配 execution log',
        '需要接入时回到 examples/INDEX.md'
    ]
)

output = '/Users/jiafan/Desktop/poc/harness-demo/docs/presentation/agent-capability-matrix.pptx'
prs.save(output)
print(f'PPTX saved: {output}')
